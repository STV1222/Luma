from __future__ import annotations
import os
from typing import Optional


TEXT_EXTS = {
    ".txt", ".md", ".py", ".json", ".csv", ".log", ".yaml", ".yml", ".ini", ".cfg", ".toml",
    ".js", ".ts", ".tsx", ".css", ".html", ".xml",
    # Common code files
    ".c", ".cpp", ".h", ".hpp", ".ino", ".java", ".kt", ".swift", ".rs", ".go", ".rb",
    ".sh", ".bash", ".zsh", ".ps1", ".bat", ".m", ".mm"
}


def _read_text_file(path: str, max_chars: int = 100_000) -> Optional[str]:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            data = f.read(max_chars)
        return data
    except Exception:
        return None


def _read_pdf(path: str, max_pages: int = 20) -> Optional[str]:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        return None
    try:
        reader = PdfReader(path)
        out: list[str] = []
        for i, page in enumerate(reader.pages):
            if i >= max_pages:
                break
            try:
                out.append(page.extract_text() or "")
            except Exception:
                continue
        text = "\n\n".join(out).strip()
        return text or None
    except Exception:
        return None


def _read_docx(path: str, max_paragraphs: int = 500) -> Optional[str]:
    try:
        import docx  # type: ignore
    except Exception:
        return None
    try:
        doc = docx.Document(path)
        paras = []
        for i, p in enumerate(doc.paragraphs):
            if i >= max_paragraphs:
                break
            paras.append(p.text)
        text = "\n".join(paras).strip()
        return text or None
    except Exception:
        return None


def _read_pptx(path: str, max_slides: int = 200, max_chars: int = 120_000) -> Optional[str]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return None
    try:
        prs = Presentation(path)
        parts: list[str] = []
        for si, slide in enumerate(prs.slides):
            if si >= max_slides:
                break
            slide_text: list[str] = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        slide_text.append(shape.text)
                except Exception:
                    continue
            if slide_text:
                parts.append("\n".join(slide_text))
            if sum(len(p) for p in parts) > max_chars:
                break
        text = "\n\n".join(parts).strip()
        return text or None
    except Exception:
        return None


def extract_text_from_file(path: str) -> Optional[str]:
    try:
        ext = os.path.splitext(path)[1].lower()
        if ext in TEXT_EXTS:
            return _read_text_file(path)
        if ext == ".pdf":
            return _read_pdf(path)
        if ext in {".docx"}:
            return _read_docx(path)
        if ext in {".pptx"}:
            return _read_pptx(path)
        # Fallback: attempt to read small non-binary files as text
        try:
            with open(path, "rb") as f:
                chunk = f.read(2048)
            if b"\x00" in chunk:
                return None  # likely binary
        except Exception:
            return None
        return _read_text_file(path)
    except Exception:
        return None



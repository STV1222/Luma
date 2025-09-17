from __future__ import annotations
import os
import io
import json
import time
import hashlib
import traceback
from dataclasses import dataclass
from datetime import datetime
from typing import Iterable, List, Tuple, Optional, Dict

# External deps expected in requirements.txt
try:
    import faiss  # type: ignore
    HAVE_FAISS = True
except Exception:
    faiss = None  # type: ignore
    HAVE_FAISS = False

try:
    from sentence_transformers import SentenceTransformer  # type: ignore
    HAVE_ST = True
except Exception:
    SentenceTransformer = None  # type: ignore
    HAVE_ST = False


RAG_HOME = os.path.expanduser("~/.luma/rag_db")
FAISS_PATH = os.path.join(RAG_HOME, "faiss.index")
META_PATH = os.path.join(RAG_HOME, "meta.jsonl")


SUPPORTED_EXTS = {
    ".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown", ".html", ".htm",
}


def ensure_dirs() -> None:
    if not os.path.isdir(RAG_HOME):
        os.makedirs(RAG_HOME, exist_ok=True)


def sha1_of_text(text: str, path: str, page: Optional[int]) -> str:
    normalized = " ".join(text.split())
    h = hashlib.sha1()
    h.update((normalized + "|" + path + "|" + str(page if page is not None else "-")).encode("utf-8", "ignore"))
    return h.hexdigest()


def load_text_from_file(path: str) -> Tuple[str, Optional[List[Tuple[int, str]]]]:
    """Return (full_text, paged_items) where paged_items is list of (page_num, page_text) if applicable.
    For PDFs/PPTX we return page-level items; for others, None.
    """
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader  # type: ignore
            reader = PdfReader(path)
            pages: List[Tuple[int, str]] = []
            for i, p in enumerate(reader.pages):
                try:
                    pages.append((i + 1, p.extract_text() or ""))
                except Exception:
                    pages.append((i + 1, ""))
            return ("\n\n".join(t for _, t in pages), pages)
        if ext == ".docx":
            import docx  # type: ignore
            doc = docx.Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)
            return (text, None)
        if ext == ".pptx":
            from pptx import Presentation  # type: ignore
            prs = Presentation(path)
            pages: List[Tuple[int, str]] = []
            for i, slide in enumerate(prs.slides):
                texts: List[str] = []
                try:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(str(shape.text))
                except Exception:
                    pass
                pages.append((i + 1, "\n".join(texts)))
            return ("\n\n".join(t for _, t in pages), pages)
        if ext in {".txt", ".md", ".markdown"}:
            # Try chardet to guess encoding
            try:
                import chardet  # type: ignore
                with open(path, "rb") as f:
                    raw = f.read()
                enc = chardet.detect(raw).get("encoding") or "utf-8"
                text = raw.decode(enc, errors="ignore")
            except Exception:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            return (text, None)
        if ext in {".html", ".htm"}:
            try:
                import chardet  # type: ignore
                with open(path, "rb") as f:
                    raw = f.read()
                enc = chardet.detect(raw).get("encoding") or "utf-8"
                html = raw.decode(enc, errors="ignore")
            except Exception:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    html = f.read()
            # dumb HTML strip
            import re
            text = re.sub(r"<[^>]+>", " ", html)
            return (text, None)
        # Fallback
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return (f.read(), None)
    except Exception:
        return ("", None)


def iter_sliding_windows(text: str, max_chars: int = 1200, overlap: int = 200) -> Iterable[str]:
    text = text.strip()
    if not text:
        return []
    # Try to keep paragraphs intact: split by double newline first.
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [text]

    chunks: List[str] = []
    buf = ""
    for p in paragraphs:
        if not buf:
            buf = p
        elif len(buf) + 2 + len(p) <= max_chars:
            buf = f"{buf}\n\n{p}"
        else:
            chunks.append(buf)
            # Sliding window overlap
            buf_tail = buf[-overlap:] if overlap > 0 else ""
            buf = (buf_tail + "\n\n" + p).strip()
            if len(buf) > max_chars:
                # Hard wrap if single paragraph too large
                for i in range(0, len(buf), max_chars - overlap):
                    sub = buf[i:i + (max_chars - overlap)]
                    if sub:
                        chunks.append(sub)
                buf = ""
    if buf:
        chunks.append(buf)
    return chunks


@dataclass
class MetaEntry:
    id: int
    path: str
    folder: str
    mtime_iso: str
    page: Optional[int]
    text_hash: str
    text: str
    deleted: bool = False


class RAGIndex:
    """Thin wrapper around FAISS IndexFlatIP plus JSONL sidecar metadata.

    - We store normalized vectors so inner product == cosine similarity.
    - We append-only for simplicity. Deletions are soft (meta.deleted = True).
    """

    def __init__(self, dim: int = 384) -> None:
        self.dim = dim
        self.index = None
        self.model = None
        self.size: int = 0  # number of vectors currently in the FAISS index

    def _lazy_model(self):
        if not HAVE_ST:
            raise RuntimeError("sentence-transformers not installed; install to enable RAG")
        if self.model is None:
            # Small, fast, widely available
            self.model = SentenceTransformer("all-MiniLM-L6-v2")
        return self.model

    def _lazy_index(self):
        if not HAVE_FAISS:
            raise RuntimeError("faiss not installed; install faiss-cpu to enable RAG")
        if self.index is None:
            if os.path.isfile(FAISS_PATH):
                try:
                    self.index = faiss.read_index(FAISS_PATH)
                    self.size = self.index.ntotal
                    return self.index
                except Exception:
                    # Corrupt index; rebuild
                    self.index = faiss.IndexFlatIP(self.dim)
                    self.size = 0
                    return self.index
            self.index = faiss.IndexFlatIP(self.dim)
            self.size = 0
        return self.index

    def _normalize(self, X):
        import numpy as np  # type: ignore
        norms = np.linalg.norm(X, axis=1, keepdims=True) + 1e-12
        return X / norms

    def _embed(self, texts: List[str]):
        import numpy as np  # type: ignore
        model = self._lazy_model()
        vecs = model.encode(texts, batch_size=64, show_progress_bar=False, convert_to_numpy=True)
        vecs = self._normalize(vecs.astype("float32"))
        return vecs

    def _append_meta(self, metas: List[MetaEntry]) -> None:
        ensure_dirs()
        temp_path = META_PATH + ".tmp"
        with open(temp_path, "a", encoding="utf-8") as w:
            for m in metas:
                w.write(json.dumps(m.__dict__, ensure_ascii=False) + "\n")
        # Atomic append: concatenate tmp to real and remove tmp
        with open(temp_path, "r", encoding="utf-8") as r, open(META_PATH, "a", encoding="utf-8") as out:
            out.write(r.read())
        os.remove(temp_path)

    def _soft_delete_path(self, path: str) -> int:
        """Mark existing entries with this path as deleted in-place by rewriting JSONL.
        Returns number of entries marked.
        """
        if not os.path.isfile(META_PATH):
            return 0
        temp = META_PATH + ".rewrite.tmp"
        changed = 0
        with open(META_PATH, "r", encoding="utf-8") as r, open(temp, "w", encoding="utf-8") as w:
            for line in r:
                try:
                    obj = json.loads(line)
                except Exception:
                    continue
                if obj.get("path") == path and not obj.get("deleted", False):
                    obj["deleted"] = True
                    changed += 1
                w.write(json.dumps(obj, ensure_ascii=False) + "\n")
        os.replace(temp, META_PATH)
        return changed

    def index_file(self, path: str) -> Dict[str, int]:
        """Index a single file. Replaces prior entries for this path via soft-delete.

        Returns a summary dict with counts.
        """
        ensure_dirs()
        ext = os.path.splitext(path)[1].lower()
        if ext not in SUPPORTED_EXTS:
            return {"added": 0, "deleted": 0}
        try:
            st = os.stat(path)
            mtime_iso = datetime.fromtimestamp(st.st_mtime).isoformat()
        except Exception:
            return {"added": 0, "deleted": 0}

        deleted = self._soft_delete_path(path)

        full_text, paged = load_text_from_file(path)
        if not full_text.strip():
            return {"added": 0, "deleted": deleted}

        # Build chunks. If paged exists, chunk per page for better citation.
        entries: List[Tuple[Optional[int], str]] = []
        if paged:
            for page_num, page_text in paged:
                for chunk in iter_sliding_windows(page_text):
                    if chunk.strip():
                        entries.append((page_num, chunk))
        else:
            for chunk in iter_sliding_windows(full_text):
                if chunk.strip():
                    entries.append((None, chunk))

        if not entries:
            return {"added": 0, "deleted": deleted}

        # Dedup by text_hash
        metas: List[MetaEntry] = []
        texts: List[str] = []
        seen: set[str] = set()
        folder = os.path.basename(os.path.dirname(path))
        for page, chunk in entries:
            h = sha1_of_text(chunk, path, page)
            if h in seen:
                continue
            seen.add(h)
            metas.append(MetaEntry(id=self.size + len(texts), path=path, folder=folder, mtime_iso=mtime_iso,
                                   page=page, text_hash=h, text=chunk, deleted=False))
            texts.append(chunk)

        if not texts:
            return {"added": 0, "deleted": deleted}

        vecs = self._embed(texts)
        index = self._lazy_index()
        index.add(vecs)
        self.size = index.ntotal

        # Persist index
        ensure_dirs()
        temp_index_path = FAISS_PATH + ".tmp"
        if HAVE_FAISS:
            faiss.write_index(index, temp_index_path)
            os.replace(temp_index_path, FAISS_PATH)

        # Append meta
        self._append_meta(metas)

        return {"added": len(texts), "deleted": deleted}

    def index_folders(self, folders: List[str], excludes: Optional[List[str]] = None) -> Dict[str, int]:
        """Recursively index supported files in given folders, respecting excludes.
        Returns counts.
        """
        excludes = excludes or []
        added = 0
        deleted = 0
        for root in folders:
            if not os.path.isdir(root):
                continue
            for dirpath, dirnames, filenames in os.walk(root):
                # prune excludes and hidden dirs
                dirnames[:] = [d for d in dirnames if not d.startswith('.') and not any(x in d for x in excludes)]
                for fn in filenames:
                    if fn.startswith('.'):
                        continue
                    path = os.path.join(dirpath, fn)
                    ext = os.path.splitext(path)[1].lower()
                    if ext not in SUPPORTED_EXTS:
                        continue
                    try:
                        res = self.index_file(path)
                        added += res.get("added", 0)
                        deleted += res.get("deleted", 0)
                    except Exception:
                        # Keep indexing even if a file fails
                        traceback.print_exc()
                        continue
        return {"added": added, "deleted": deleted}


def read_meta_lines() -> Iterable[Dict[str, object]]:
    if not os.path.isfile(META_PATH):
        return []
    def _iter():
        with open(META_PATH, "r", encoding="utf-8") as r:
            for line in r:
                try:
                    yield json.loads(line)
                except Exception:
                    continue
    return _iter()


